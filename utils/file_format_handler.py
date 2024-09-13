#This file is for handling different types of file formats such as pdf,ppt,txt,docx etc
from PyPDF2 import PdfReader
from docx import Document 
from pptx import Presentation
import pandas as pd
import base64
import streamlit as st

def get_pdf_text(file):  # get all text from pdf file
    pdf_reader = PdfReader(file)
    pages = []
    for i, page in enumerate(pdf_reader.pages):
        text = page.extract_text()
        if text:
            pages.append((text, i + 1))  # store text with page number (1-based index)
    return pages

def displayPDF(uploaded_file):
    # Check if the uploaded file is a valid PDF
    if uploaded_file is not None and uploaded_file.type == 'application/pdf':
        # Read the file-like object directly
        base64_pdf = base64.b64encode(uploaded_file.read()).decode('utf-8')

        # Embedding PDF in HTML
        pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="400" type="application/pdf"></iframe>'

        # Display the embedded PDF
        st.markdown(pdf_display, unsafe_allow_html=True)
    else:
        st.warning("Can only preview PDF file.")


def get_text(file, page_size=2000):
    file_content = file.read().decode('utf-8')  # Decode the content of the uploaded file
    chunks = [file_content[i:i+page_size] for i in range(0, len(file_content), page_size)]
    pages = [(chunk, i+1) for i, chunk in enumerate(chunks)]
    return pages

def get_word_text(file):
    document = Document(file)
    full_text = []
    for para in document.paragraphs:
        full_text.append(para.text)
    content = "\n".join(full_text)
    return [(content, 1)]

def get_ppt_text(file):
    presentation = Presentation(file)
    slides = []
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides.append(("\n".join(slide_text), i + 1))  # store text with slide number (1-based index)
    return slides

def get_excel_text(file):
    df = pd.read_excel(file)
    rows = []
    for index, row in df.iterrows():
        row_text = row.to_string(index=False)
        rows.append((row_text, index + 1))  # store text with row number (1-based index)
    return rows

def get_csv_text(file):
    df = pd.read_csv(file)
    rows = []
    for index, row in df.iterrows():
        row_text = row.to_string(index=False)
        rows.append((row_text, index + 1))  # store text with row number (1-based index)
    return rows


